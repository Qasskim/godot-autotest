"""
리듬게임 테스트 자동화 매뉴얼 PPTX 생성
출력: D:\Automation\RhythmGame_AutoTest\테스트자동화_매뉴얼.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy, os

OUT = r"D:\Automation\RhythmGame_AutoTest\테스트자동화_매뉴얼.pptx"
IMG_DIR = r"D:\Automation\RhythmGame_AutoTest\refs"

# ── 색상 ──────────────────────────────────────────────
BG      = RGBColor(0x1A, 0x1A, 0x2E)   # 짙은 남색
BG2     = RGBColor(0x16, 0x21, 0x3E)   # 카드 배경
CYAN    = RGBColor(0x00, 0xD4, 0xFF)   # 포인트
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0xB0, 0xB8, 0xC8)
GREEN   = RGBColor(0x2E, 0xCC, 0x71)
YELLOW  = RGBColor(0xF3, 0x9C, 0x12)
RED     = RGBColor(0xE7, 0x4C, 0x3C)
DARKBG  = RGBColor(0x0F, 0x0F, 0x23)

W = Inches(13.33)
H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color: RGBColor, alpha=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Malgun Gothic"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txb


def multi_line_txt(slide, lines, x, y, w, h, size=14, color=WHITE,
                   font="Malgun Gothic", line_spacing=None):
    """(text, bold, color) 튜플 리스트로 멀티라인 텍스트박스 생성"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, col = item, False, color
        else:
            text, bold, col = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if line_spacing:
            from pptx.util import Pt as UPt
            from pptx.oxml.ns import qn
            from lxml import etree
            pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
            spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
            spcPts.set('val', str(int(line_spacing * 100)))
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = col
        run.font.name = font
    return txb


def add_image(slide, path, x, y, w, h=None):
    if h:
        slide.shapes.add_picture(path, x, y, w, h)
    else:
        slide.shapes.add_picture(path, x, y, w)


def card(slide, x, y, w, h, color=BG2, radius=False):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = CYAN
    shape.line.width = Pt(0.5)
    return shape


def cyan_label(slide, text, x, y, w=None):
    w = w or Inches(3)
    r = rect(slide, x, y, w, Inches(0.38), CYAN)
    txt(slide, text, x + Inches(0.1), y + Pt(3), w, Inches(0.38),
        size=13, bold=True, color=DARKBG)


def step_circle(slide, num, x, y):
    r = rect(slide, x, y, Inches(0.5), Inches(0.5), CYAN)
    txt(slide, str(num), x, y + Pt(2), Inches(0.5), Inches(0.5),
        size=14, bold=True, color=DARKBG, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# Slide 1 — 타이틀
# ════════════════════════════════════════════════════════
def slide_title(prs):
    sl = blank_slide(prs)
    fill_bg(sl, DARKBG)

    # 배경 사각형 (왼쪽 강조 블록)
    rect(sl, 0, 0, Inches(0.5), H, CYAN)
    rect(sl, Inches(0.5), 0, Inches(12.83), H, BG)

    # 좌측 세로 텍스트 라인 장식
    rect(sl, Inches(1.2), Inches(1.5), Pt(2), Inches(4.5), CYAN)

    # 타이틀
    txt(sl, "리듬게임 테스트 자동화", Inches(1.5), Inches(1.6),
        Inches(9), Inches(1.2), size=44, bold=True, color=WHITE)
    txt(sl, "매뉴얼 v1.0", Inches(1.5), Inches(2.7),
        Inches(9), Inches(0.7), size=28, bold=False, color=CYAN)

    # 서브 정보
    info_lines = [
        ("프로젝트", True, CYAN),
        ("  Godot 4.4 Android 리듬게임 (16레인 듀얼 윙)", False, GRAY),
        ("", False, WHITE),
        ("테스트 방식", True, CYAN),
        ("  Godot 내장 HTTP TestServer (포트 5000) + Appium", False, GRAY),
        ("", False, WHITE),
        ("대상 기기", True, CYAN),
        ("  실기기  Samsung Galaxy Note 20 (R3CR501XHAJ)", False, GRAY),
        ("  AVD     RhythmGame_FHD (emulator-5554, 2280×1080)", False, GRAY),
        ("", False, WHITE),
        ("TC 수", True, CYAN),
        ("  12건 전체 PASS", False, GREEN),
    ]
    multi_line_txt(sl, info_lines, Inches(1.5), Inches(3.5),
                   Inches(7), Inches(3.5), size=14)

    # 우측 게임 스크린샷
    add_image(sl, os.path.join(IMG_DIR, "04_ingame.png"),
              Inches(9), Inches(1.2), Inches(4), Inches(2.25))
    add_image(sl, os.path.join(IMG_DIR, "06_result.png"),
              Inches(9), Inches(3.55), Inches(4), Inches(2.25))

    txt(sl, "Qasskim / github.com/Qasskim/godot-autotest",
        Inches(1.5), Inches(6.9), Inches(8), Inches(0.5),
        size=11, color=GRAY)


# ════════════════════════════════════════════════════════
# Slide 2 — 아키텍처 개요
# ════════════════════════════════════════════════════════
def slide_arch(prs):
    sl = blank_slide(prs)
    fill_bg(sl, BG)
    rect(sl, 0, 0, Inches(13.33), Inches(0.9), DARKBG)
    txt(sl, "아키텍처 개요", Inches(0.5), Inches(0.1),
        Inches(8), Inches(0.8), size=28, bold=True, color=CYAN)

    # 왜 HTTP 방식인가?
    cyan_label(sl, "기존 방식의 한계", Inches(0.5), Inches(1.05))
    rect(sl, Inches(0.5), Inches(1.5), Inches(5.5), Inches(1.8), RGBColor(0x2C,0x10,0x10))
    multi_line_txt(sl, [
        ("UIAutomator2", True, RED),
        ("  Godot GL 캔버스 → 전체가 단 하나의 View", False, GRAY),
        ("  개별 버튼 탐지 불가, accessibility_name 무효", False, GRAY),
        ("이미지 매칭", True, RED),
        ("  토글 ON/OFF 상태 판별 불가 (신뢰성 없음)", False, GRAY),
    ], Inches(0.6), Inches(1.55), Inches(5.3), Inches(1.7), size=13)

    # 화살표
    txt(sl, "→", Inches(6.1), Inches(2.0), Inches(0.6), Inches(0.6),
        size=28, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

    # 채택된 방식
    cyan_label(sl, "채택: Godot HTTP TestServer", Inches(6.8), Inches(1.05))
    rect(sl, Inches(6.8), Inches(1.5), Inches(6.0), Inches(1.8), RGBColor(0x0A,0x2A,0x1A))
    multi_line_txt(sl, [
        ("Godot 앱 내부에 경량 HTTP 서버 내장 (포트 5000)", True, GREEN),
        ("  debug 빌드에서만 활성화 (OS.is_debug_build())", False, GRAY),
        ("adb forward tcp:5000 → PC에서 localhost:5000 접근", False, GRAY),
        ("Appium은 앱 실행/종료만 담당", False, GRAY),
        ("나머지 모든 제어는 HTTP API로 처리", False, GRAY),
    ], Inches(6.9), Inches(1.55), Inches(5.8), Inches(1.7), size=13)

    # 구성도
    cyan_label(sl, "구성도", Inches(0.5), Inches(3.5))

    # 박스들
    boxes = [
        (Inches(0.5),  Inches(4.0), "PC\nPython 테스트"),
        (Inches(3.8),  Inches(4.0), "Appium Server\nlocalhost:4723"),
        (Inches(7.1),  Inches(4.0), "Android 기기\n(실기기/AVD)"),
        (Inches(10.4), Inches(4.0), "Godot\nTestServer :5000"),
    ]
    for bx, by, label in boxes:
        card(sl, bx, by, Inches(3.0), Inches(1.2))
        txt(sl, label, bx + Inches(0.15), by + Inches(0.2),
            Inches(2.7), Inches(0.9), size=13, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)

    # 화살표 연결
    arrows = [
        (Inches(3.55), Inches(4.6), "앱 실행/종료"),
        (Inches(6.85), Inches(4.6), "adb forward"),
        (Inches(10.15), Inches(4.6), "HTTP GET/POST"),
    ]
    for ax, ay, label in arrows:
        txt(sl, "→", ax, ay - Inches(0.05), Inches(0.35), Inches(0.5),
            size=20, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
        txt(sl, label, ax - Inches(0.1), ay + Inches(0.3), Inches(0.6), Inches(0.4),
            size=10, color=GRAY, align=PP_ALIGN.CENTER)

    # API 목록
    cyan_label(sl, "HTTP API 엔드포인트", Inches(0.5), Inches(5.5))
    api_items = [
        "GET  /ping              앱 생존 확인",
        "GET  /scene             현재 씬 이름",
        "GET  /autoplay          오토플레이 상태",
        "GET  /result            결과 데이터 (score, grade, wow, miss...)",
        "POST /autoplay/on|off   오토플레이 제어",
        "POST /navigate/...      씬 전환 트리거",
        "POST /tap/...           버튼 탭 트리거",
    ]
    multi_line_txt(sl,
        [(it, False, GRAY) for it in api_items],
        Inches(0.5), Inches(5.9), Inches(12.5), Inches(1.5),
        size=12, font="Consolas")


# ════════════════════════════════════════════════════════
# Slide 3 — 사전 준비
# ════════════════════════════════════════════════════════
def slide_setup(prs):
    sl = blank_slide(prs)
    fill_bg(sl, BG)
    rect(sl, 0, 0, Inches(13.33), Inches(0.9), DARKBG)
    txt(sl, "사전 준비", Inches(0.5), Inches(0.1),
        Inches(8), Inches(0.8), size=28, bold=True, color=CYAN)

    col1_x = Inches(0.5)
    col2_x = Inches(6.8)
    y0 = Inches(1.0)

    # ── 왼쪽: 필수 소프트웨어 ──
    cyan_label(sl, "필수 소프트웨어 설치", col1_x, y0)
    sw_items = [
        ("Python 3.10+", "pip install -r requirements.txt"),
        ("Appium 2.x", "npm install -g appium  |  appium driver install uiautomator2"),
        ("Android SDK", "adb 명령어 PATH 등록 필요"),
        ("Godot 4.4.1", r"D:\Godot_v4.4.1-stable_win64.exe\... (빌드용)"),
        ("Java 11+", "Appium / UIAutomator2 의존성"),
    ]
    for i, (name, desc) in enumerate(sw_items):
        y = y0 + Inches(0.5) + i * Inches(0.75)
        step_circle(sl, i+1, col1_x, y)
        txt(sl, name, col1_x + Inches(0.65), y,
            Inches(5.5), Inches(0.35), size=13, bold=True, color=WHITE)
        txt(sl, desc, col1_x + Inches(0.65), y + Inches(0.3),
            Inches(5.5), Inches(0.35), size=11, color=GRAY, font="Consolas")

    # ── 오른쪽: 경로 / 설정 ──
    cyan_label(sl, "경로 설정 (main.py)", col2_x, y0)
    path_lines = [
        ("ADB", True, CYAN),
        (r"  %LOCALAPPDATA%\Android\Sdk\platform-tools\adb.exe", False, GRAY),
        ("", False, WHITE),
        ("GODOT", True, CYAN),
        (r"  D:\Godot_v4.4.1-stable_win64.exe\..._console.exe", False, GRAY),
        ("", False, WHITE),
        ("PROJECT", True, CYAN),
        (r"  D:\Godot\RhythmGame_New\Godot", False, GRAY),
        ("", False, WHITE),
        ("APK 출력", True, CYAN),
        (r"  D:\Automation\apk\RhythmGame_debug.apk", False, GRAY),
    ]
    multi_line_txt(sl, path_lines, col2_x, y0 + Inches(0.5),
                   Inches(6.2), Inches(3.5), size=12, font="Consolas")

    # ── 하단: 기기 연결 확인 ──
    cyan_label(sl, "기기 연결 확인", col1_x, Inches(5.1))
    check_cmds = [
        "adb devices                        # 연결된 기기 목록 확인",
        "adb -s R3CR501XHAJ get-model       # 실기기 연결 확인",
        "emulator -list-avds                # AVD 목록 확인",
        "emulator -avd RhythmGame_FHD       # AVD 실행 (별도 터미널)",
    ]
    multi_line_txt(sl,
        [(c, False, GRAY) for c in check_cmds],
        col1_x, Inches(5.55), Inches(12.5), Inches(1.8),
        size=12, font="Consolas")

    # requirements
    cyan_label(sl, "requirements.txt", col2_x, Inches(4.1))
    reqs = [
        "Appium-Python-Client>=3.0.0",
        "openpyxl>=3.1.0",
        "Pillow>=10.0.0",
        "pytesseract>=0.3.10",
        "psutil>=5.9.0",
        "requests",
    ]
    multi_line_txt(sl,
        [(r, False, GRAY) for r in reqs],
        col2_x, Inches(4.55), Inches(5.5), Inches(1.5),
        size=12, font="Consolas")


# ════════════════════════════════════════════════════════
# Slide 4 — 실행 전 준비 (초보자용 단계별)
# ════════════════════════════════════════════════════════
def slide_run_prereq(prs):
    sl = blank_slide(prs)
    fill_bg(sl, BG)
    rect(sl, 0, 0, Inches(13.33), Inches(0.9), DARKBG)
    txt(sl, "실행 전 준비  —  매번 실행 전 확인하는 순서", Inches(0.5), Inches(0.1),
        Inches(12), Inches(0.8), size=26, bold=True, color=CYAN)

    steps = [
        {
            "num": "STEP 1",
            "title": "터미널 A  —  Appium 서버 시작",
            "color": CYAN,
            "cmd": "appium",
            "note": "Appium은 Python 테스트 코드와 Android 기기 사이의 중간 서버입니다.\n"
                    "테스트를 실행하기 전에 반드시 먼저 켜 두어야 합니다.\n"
                    "실행 후 'Appium REST http interface listener started on 0.0.0.0:4723' 메시지가 보이면 준비 완료.",
            "warn": None,
        },
        {
            "num": "STEP 2",
            "title": "기기 연결 확인  (실기기 사용 시)",
            "color": GREEN,
            "cmd": "adb devices",
            "note": "USB 케이블로 연결 후 위 명령어를 실행합니다.\n"
                    "결과에 'R3CR501XHAJ   device' 가 보이면 연결 성공.\n"
                    "'unauthorized' 가 보이면 기기 화면에서 'USB 디버깅 허용' 팝업을 수락하세요.",
            "warn": "기기 설정 > 개발자 옵션 > USB 디버깅 이 켜져 있어야 합니다.",
        },
        {
            "num": "STEP 2",
            "title": "AVD 실행  (에뮬레이터 사용 시)  —  터미널 B",
            "color": YELLOW,
            "cmd": "emulator -avd RhythmGame_FHD",
            "note": "별도 터미널 창을 열어서 실행합니다. 에뮬레이터 화면이 뜨고\n"
                    "Android 홈 화면이 보일 때까지 기다리세요 (1~2분 소요).\n"
                    "이후 'adb devices' 로 'emulator-5554  device' 가 뜨면 준비 완료.",
            "warn": None,
        },
        {
            "num": "STEP 3",
            "title": "터미널 C  —  테스트 실행 (다음 슬라이드 참고)",
            "color": RGBColor(0xAA, 0xAA, 0xCC),
            "cmd": "cd D:\\Automation\\RhythmGame_AutoTest",
            "note": "테스트 폴더로 이동한 뒤 다음 슬라이드의 명령어를 실행합니다.\n"
                    "Appium 서버(터미널 A)와 기기/에뮬레이터는 계속 켜 두어야 합니다.",
            "warn": None,
        },
    ]

    col_l = Inches(0.45)
    col_r = Inches(6.85)
    positions = [
        (col_l, Inches(1.0)),
        (col_l, Inches(3.55)),
        (col_r, Inches(1.0)),
        (col_r, Inches(3.55)),
    ]

    for i, (step, (sx, sy)) in enumerate(zip(steps, positions)):
        card_w = Inches(6.2)
        card_h = Inches(2.35)
        card(sl, sx, sy, card_w, card_h)

        # 번호 뱃지
        num_color = step["color"]
        rect(sl, sx, sy, Inches(1.0), Inches(0.38), num_color)
        txt(sl, step["num"], sx, sy, Inches(1.0), Inches(0.38),
            size=11, bold=True, color=DARKBG, align=PP_ALIGN.CENTER)

        # 제목
        txt(sl, step["title"], sx + Inches(1.1), sy + Pt(3),
            card_w - Inches(1.2), Inches(0.38),
            size=13, bold=True, color=WHITE)

        # 명령어 박스
        cmd_y = sy + Inches(0.45)
        rect(sl, sx + Inches(0.12), cmd_y, card_w - Inches(0.25), Inches(0.38),
             DARKBG)
        txt(sl, step["cmd"], sx + Inches(0.22), cmd_y + Pt(3),
            card_w - Inches(0.35), Inches(0.35),
            size=13, bold=True, color=num_color, font="Consolas")

        # 설명
        note_y = cmd_y + Inches(0.45)
        txt(sl, step["note"], sx + Inches(0.15), note_y,
            card_w - Inches(0.25), Inches(1.1),
            size=11, color=GRAY)

        # 경고
        if step["warn"]:
            warn_y = note_y + Inches(0.85)
            rect(sl, sx + Inches(0.12), warn_y, card_w - Inches(0.25), Inches(0.32),
                 RGBColor(0x4A, 0x2A, 0x00))
            txt(sl, "⚠  " + step["warn"], sx + Inches(0.2), warn_y + Pt(2),
                card_w - Inches(0.35), Inches(0.3),
                size=10, color=YELLOW)

    # 실기기/AVD 선택 안내
    txt(sl, "※  실기기와 AVD 중 하나만 선택하면 됩니다. STEP 2는 사용하는 방식에 맞게 진행하세요.",
        Inches(0.45), Inches(6.95), Inches(12.4), Inches(0.45),
        size=12, color=GRAY, italic=True)


# ════════════════════════════════════════════════════════
# Slide 5 — 실행 명령어
# ════════════════════════════════════════════════════════
def slide_run(prs):
    sl = blank_slide(prs)
    fill_bg(sl, BG)
    rect(sl, 0, 0, Inches(13.33), Inches(0.9), DARKBG)
    txt(sl, "테스트 실행  —  명령어 가이드", Inches(0.5), Inches(0.1),
        Inches(10), Inches(0.8), size=26, bold=True, color=CYAN)

    # ── 실행 명령어 4종 ──
    cmds = [
        ("실기기 (기본)",
         "python main.py",
         "실기기 R3CR501XHAJ 에서 실행합니다.\n"
         "APK를 자동으로 빌드(Godot) → 설치(adb) → 포트 포워딩 → 테스트 순서로 진행됩니다.\n"
         "처음 빌드는 1~3분 소요됩니다."),
        ("AVD (에뮬레이터)",
         "python main.py emulator-5554",
         "에뮬레이터에서 실행합니다.\n"
         "먼저 에뮬레이터가 완전히 부팅되어 있어야 합니다 (이전 슬라이드 STEP 2 참고)."),
        ("복수 기기 순차 실행",
         "python main.py R3CR501XHAJ emulator-5554",
         "실기기와 AVD 모두에서 순차적으로 실행합니다.\n"
         "결과 xlsx 파일이 기기별로 각각 생성됩니다."),
        ("빌드 생략 (빠른 재실행)",
         "python main.py --skip-build emulator-5554",
         "이미 APK가 있을 때 빌드 단계를 건너뜁니다.\n"
         "코드만 수정하고 APK는 그대로 쓸 때 사용 (시간 단축)."),
    ]

    for i, (label, cmd, desc) in enumerate(cmds):
        y = Inches(1.0) + i * Inches(1.5)
        card(sl, Inches(0.45), y, Inches(12.4), Inches(1.38))

        # 레이블
        txt(sl, label, Inches(0.65), y + Inches(0.07),
            Inches(4), Inches(0.32), size=12, bold=True, color=CYAN)

        # 명령어
        rect(sl, Inches(0.45), y + Inches(0.38), Inches(12.4), Inches(0.42), DARKBG)
        txt(sl, ">  " + cmd, Inches(0.6), y + Inches(0.4),
            Inches(12.0), Inches(0.38), size=14, bold=True,
            color=GREEN, font="Consolas")

        # 설명
        txt(sl, desc, Inches(0.65), y + Inches(0.85),
            Inches(12.0), Inches(0.5), size=11, color=GRAY)

    # ── 실행 후 일어나는 일 (내부 흐름) ──
    flow_y = Inches(7.05)
    txt(sl, "python main.py 실행 시 내부 순서 →", Inches(0.45), flow_y - Inches(0.28),
        Inches(4), Inches(0.3), size=11, color=GRAY)

    steps = [
        ("① Godot\nAPK 빌드", CYAN),
        ("② adb\ninstall -r", CYAN),
        ("③ adb forward\ntcp:5000", CYAN),
        ("④ Appium\n드라이버 연결", CYAN),
        ("⑤ TC 01~12\n순차 실행", CYAN),
        ("⑥ test_results\n.xlsx 저장", GREEN),
    ]
    step_w = Inches(1.95)
    for i, (label, col) in enumerate(steps):
        bx = Inches(0.45) + i * (step_w + Inches(0.2))
        rect(sl, bx, flow_y, step_w, Inches(0.58), RGBColor(0x0A, 0x2A, 0x3A))
        txt(sl, label, bx + Inches(0.05), flow_y + Inches(0.04),
            step_w - Inches(0.1), Inches(0.55), size=11,
            bold=False, color=col, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            txt(sl, "→", bx + step_w + Inches(0.03), flow_y + Inches(0.1),
                Inches(0.18), Inches(0.4), size=14, bold=True,
                color=CYAN, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# Slide 5 — TC 목록
# ════════════════════════════════════════════════════════
def slide_tc(prs):
    sl = blank_slide(prs)
    fill_bg(sl, BG)
    rect(sl, 0, 0, Inches(13.33), Inches(0.9), DARKBG)
    txt(sl, "테스트 케이스 목록 (TC 12건)", Inches(0.5), Inches(0.1),
        Inches(10), Inches(0.8), size=28, bold=True, color=CYAN)
    txt(sl, "전체 PASS ✓", Inches(10.5), Inches(0.15),
        Inches(2.5), Inches(0.7), size=20, bold=True, color=GREEN,
        align=PP_ALIGN.RIGHT)

    tc_list = [
        ("TC-01", "test_autoplay",    "앱 실행 후 메인 화면에 진입했는가?",                   "MainMenu 씬 확인"),
        ("TC-02", "test_autoplay",    "PLAY 버튼 터치 후 곡 선택 화면에 진입했는가?",         "SongSelect 씬 확인"),
        ("TC-03", "test_autoplay",    "DEBUG에서 오토플레이를 활성화했는가?",                  "API auto_play=true 확인"),
        ("TC-04", "test_autoplay",    "곡 선택 후 게임 화면에 진입했는가?",                   "GameScene 씬 확인"),
        ("TC-05", "test_autoplay",    "오토플레이 완주 후 결과 화면에 진입했는가?",            "ResultScreen 씬 확인 (최대 300초)"),
        ("TC-06", "test_autoplay",    "오토플레이 결과에서 Miss가 0인가?",                    "GET /result → miss==0"),
        ("TC-07", "test_autoplay",    "오토플레이 결과 등급이 MAX인가?",                      "GET /result → grade==\"MAX\""),
        ("TC-08", "test_result",      "결과 화면에 판정 항목이 모두 표시되는가?",             "wow/great/good/oops 키 존재 확인"),
        ("TC-09", "test_result",      "결과 화면에 점수가 표시되는가?",                       "score >= 0 확인"),
        ("TC-10", "test_result",      "다시하기 버튼 동작 후 게임 화면에 진입했는가?",        "tap_retry → GameScene"),
        ("TC-11", "test_result",      "곡 선택 버튼 동작 후 곡 선택 화면에 복귀했는가?",     "tap_select_song → SongSelect"),
        ("TC-12", "test_navigation",  "곡 선택 화면에서 메인 메뉴로 복귀했는가?",            "navigate_main_menu → MainMenu"),
    ]

    module_colors = {
        "test_autoplay":   CYAN,
        "test_result":     YELLOW,
        "test_navigation": GREEN,
    }

    col_x  = [Inches(0.45), Inches(1.3), Inches(3.55), Inches(9.3)]
    header = ["번호", "모듈", "테스트 항목", "검증 방법"]
    header_y = Inches(1.0)
    rect(sl, Inches(0.4), header_y, Inches(12.5), Inches(0.38), DARKBG)
    for i, h in enumerate(header):
        txt(sl, h, col_x[i], header_y + Pt(3),
            Inches(2.2), Inches(0.35), size=12, bold=True, color=CYAN)

    for row, (tc, mod, item, verify) in enumerate(tc_list):
        y = header_y + Inches(0.42) + row * Inches(0.46)
        bg_color = RGBColor(0x1A,0x21,0x30) if row % 2 == 0 else BG
        rect(sl, Inches(0.4), y, Inches(12.5), Inches(0.44), bg_color)
        mc = module_colors.get(mod, GRAY)
        txt(sl, tc, col_x[0], y + Pt(4), Inches(0.8), Inches(0.38),
            size=12, bold=True, color=WHITE, font="Consolas")
        txt(sl, mod, col_x[1], y + Pt(4), Inches(2.1), Inches(0.38),
            size=10, color=mc, font="Consolas")
        txt(sl, item, col_x[2], y + Pt(4), Inches(5.6), Inches(0.38),
            size=12, color=WHITE)
        txt(sl, verify, col_x[3], y + Pt(4), Inches(3.8), Inches(0.38),
            size=11, color=GRAY)


# ════════════════════════════════════════════════════════
# Slide 6 — 화면 흐름 + 스크린샷
# ════════════════════════════════════════════════════════
def slide_flow(prs):
    sl = blank_slide(prs)
    fill_bg(sl, BG)
    rect(sl, 0, 0, Inches(13.33), Inches(0.9), DARKBG)
    txt(sl, "실행 흐름 & 화면 스크린샷", Inches(0.5), Inches(0.1),
        Inches(10), Inches(0.8), size=28, bold=True, color=CYAN)

    screens = [
        ("01 메인 메뉴", "01_main.png",   "TC-01: 앱 실행\n→ MainMenu 씬 확인"),
        ("02 곡 선택",   "02_songselect.png", "TC-02~04: PLAY 진입\n오토플레이 ON\n첫 곡 선택 & 시작"),
        ("03 인게임",    "04_ingame.png",  "TC-05: 오토플레이\n완주 대기 (최대 300초)"),
        ("04 결과 화면", "06_result.png",  "TC-06~09: Miss=0\nGRADE MAX 확인\n판정/점수 검증"),
    ]

    img_w = Inches(2.9)
    img_h = Inches(1.63)  # 16:9 비율

    for i, (label, img, desc) in enumerate(screens):
        x = Inches(0.5) + i * (img_w + Inches(0.55))
        # 이미지
        img_path = os.path.join(IMG_DIR, img)
        add_image(sl, img_path, x, Inches(1.0), img_w, img_h)
        # 레이블
        rect(sl, x, Inches(1.0), img_w, Inches(0.3), RGBColor(0x00,0x00,0x00))
        txt(sl, label, x + Inches(0.05), Inches(1.0),
            img_w, Inches(0.3), size=11, bold=True, color=CYAN)
        # 설명
        txt(sl, desc, x, Inches(2.7),
            img_w, Inches(0.9), size=11, color=GRAY)

    # 하단 흐름도
    cyan_label(sl, "TC 실행 흐름", Inches(0.5), Inches(3.75))

    flow_items = [
        ("앱 실행", "TC-01"),
        ("곡 선택 진입", "TC-02"),
        ("오토플레이 ON", "TC-03"),
        ("첫 곡 플레이", "TC-04"),
        ("완주 대기", "TC-05"),
        ("결과 검증", "TC-06~09"),
        ("다시하기", "TC-10"),
        ("완주 재대기", "(자동)"),
        ("곡선택 복귀", "TC-11"),
        ("메인 복귀", "TC-12"),
    ]

    box_w = Inches(1.15)
    box_h = Inches(0.75)
    start_x = Inches(0.4)
    row_y = [Inches(4.25), Inches(5.3)]

    for i, (label, tc_no) in enumerate(flow_items):
        row = i // 5
        col = i % 5
        x = start_x + col * (box_w + Inches(0.15))
        y = row_y[row]
        card(sl, x, y, box_w, box_h)
        txt(sl, label, x + Inches(0.05), y + Inches(0.04),
            box_w - Inches(0.1), Inches(0.38), size=11, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, tc_no, x + Inches(0.05), y + Inches(0.42),
            box_w - Inches(0.1), Inches(0.28), size=10,
            color=CYAN, align=PP_ALIGN.CENTER)

        # 화살표 (오른쪽 또는 아래)
        if col < 4:
            ax = x + box_w + Inches(0.02)
            ay = y + Inches(0.22)
            txt(sl, "→", ax, ay, Inches(0.14), Inches(0.35),
                size=14, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
        elif row == 0:  # 줄바꿈
            txt(sl, "↓", x + box_w // 2 - Inches(0.1), y + box_h,
                Inches(0.3), Inches(0.3), size=14, bold=True,
                color=CYAN, align=PP_ALIGN.CENTER)

    # 오른쪽: 결과 파일 정보
    cyan_label(sl, "결과 출력", Inches(6.8), Inches(3.75))
    res_lines = [
        ("test_results_<기기명>.xlsx", True, WHITE),
        ("", False, WHITE),
        ("컬럼: TC번호 / 항목 / PASS/FAIL", False, GRAY),
        ("       실패 시 스크린샷 경로 / 에러 메시지", False, GRAY),
        ("", False, WHITE),
        ("스크린샷: screenshots/<datetime>_tc<N>.png", False, GRAY),
    ]
    multi_line_txt(sl, res_lines,
                   Inches(6.8), Inches(4.25), Inches(6.0), Inches(1.5),
                   size=13, font="Consolas")


# ════════════════════════════════════════════════════════
# 생성
# ════════════════════════════════════════════════════════
prs = new_prs()
slide_title(prs)
slide_arch(prs)
slide_setup(prs)
slide_run_prereq(prs)   # NEW: 실행 전 준비 (초보자용)
slide_run(prs)
slide_tc(prs)
slide_flow(prs)
prs.save(OUT)
print(f"저장 완료: {OUT}")
